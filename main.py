import streamlit as st
from openai import OpenAI
import os
from pypdf import PdfReader
from pptx import Presentation
from docx import Document

# --- PAGE CONFIGURATION & CUSTOM CSS ---
st.set_page_config(page_title="METU-IE Summer Practice Bot", page_icon="🔴", layout="centered")

st.markdown("""
    <style>
    /* Expand the main container so the header reaches the sides */
    .block-container { 
        padding-top: 2rem; 
        max-width: 95% !important; 
    }
    
    /* Make the sidebar slightly wider to enlarge the logo */
    [data-testid="stSidebar"] {
        min-width: 330px !important;
        max-width: 330px !important;
    }
    
    /* Reduce padding around the logo and language toggle */
    [data-testid="stSidebar"] .stImage {
        margin-top: -5px;
        margin-bottom: -5px;
    }
    [data-testid="stSidebar"] .stRadio {
        margin-top: -15px;
        margin-bottom: -10px;
    }
    
    /* Make the language toggle more narrow/compact */
    div[role="radiogroup"] {
        gap: 10px !important;
    }

    .metu-topbar {
        background-color: #333333;
        color: #dddddd;
        padding: 6px 15px;
        font-size: 12px;
        border-radius: 6px 6px 0 0;
        text-align: center; /* Centered text */
        font-family: Arial, sans-serif;
    }
    
    .metu-navbar {
        background-color: #cb2c30; 
        color: white;
        padding: 12px 0;
        text-align: center;
        font-weight: 700;
        font-size: 22px;
        margin-bottom: 30px;
        border-radius: 0 0 6px 6px;
        box-shadow: 0 4px 6px rgba(0,0,0,0.1);
        letter-spacing: 0.5px;
    }
    
    .header-text {
        color: #2c3e50;
        font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
        margin-top: 15px;
        font-size: 28px;
        font-weight: 600;
        line-height: 1.2;
        text-align: center;
    }

    .stButton>button {
        width: 100%;
        text-align: left;
        background-color: transparent;
        border: 1px solid #ccc;
        color: #2c3e50;
        transition: all 0.2s;
    }
    .stButton>button:hover {
        border-color: #cb2c30;
        color: #cb2c30;
    }
    
    .sidebar-link {
        display: block;
        padding: 10px 15px;
        margin-bottom: 10px;
        background-color: #f8f9fa;
        border-left: 5px solid #cb2c30;
        color: #333 !important;
        text-decoration: none;
        border-radius: 4px;
        font-weight: 500;
        transition: all 0.2s ease-in-out;
        box-shadow: 0 1px 3px rgba(0,0,0,0.05);
    }
    .sidebar-link:hover {
        background-color: #e9ecef;
        transform: translateX(3px);
    }
    </style>
""", unsafe_allow_html=True)

# --- BILINGUAL DICTIONARY ---
lang_dict = {
    "ENG": {
        "topbar": "IE | METU",
        "title": "Industrial Engineering<br>Summer Practice",
        "subtitle": "Intelligent Chatbot Assistant for IE300/IE400",
        "faq_header": "❓ Quick Questions",
        "links_header": "🔗 Quick Links",
        "contact_hdr": "📞 Contact Information",
        "link_curr": "IE Curriculum",
        "link_sp": "IE SP Webpage",
        "link_main": "IE Main Page",
        "cat1": "General & Paperwork",
        "cat2": "IE 300 & IE 400 Specifics",
        "q1_lbl": "How to arrange an organization?",
        "q1_prmpt": "How can I arrange a summer practice organization?",
        "q2_lbl": "What is the required paperwork?",
        "q2_prmpt": "What is the paperwork needed to get officially registered to summer practice?",
        "q3_lbl": "Which departments should I visit?",
        "q3_prmpt": "Which departments of the organization should I visit during my summer practice?",
        "q4_lbl": "Are banks suitable for IE 400?",
        "q4_prmpt": "What sort of branches or divisions in a bank are suitable for an IE 400 practice organization?",
        "q5_lbl": "Manufacturing rules in service firms?",
        "q5_prmpt": "How should a student answer questions in manufacturing oriented areas if the practice is done in a service organization?",
        "q6_lbl": "Sufficient problem definition (IE400)?",
        "q6_prmpt": "What constitutes a 'sufficient' IE problem definition in IE 400 reports?",
        "q7_lbl": "Can I do a full-time project instead?",
        "q7_prmpt": "If I am planning to work on a project full-time at the summer practice organization, is it possible for me to present the project work?",
        "q8_lbl": "Rules for multiple students at one firm?",
        "q8_prmpt": "What sort of reporting is acceptable if several students practice in the same organization?",
        "placeholder": "E.g., What are the requirements for IE 300?",
        "thinking": "Consulting official guidelines...",
        "system_instruction_lang": "You must respond entirely in English."
    },
    "TR": {
        "topbar": "EM | ODTÜ",
        "title": "Endüstri Mühendisliği<br>Yaz Stajı",
        "subtitle": "IE300/IE400 için Akıllı Sohbet Robotu Asistanı",
        "faq_header": "❓ Hızlı Sorular",
        "links_header": "🔗 Hızlı Bağlantılar",
        "contact_hdr": "📞 İletişim Bilgileri",
        "link_curr": "EM Müfredatı",
        "link_sp": "EM Yaz Stajı Sayfası",
        "link_main": "EM Ana Sayfası",
        "cat1": "Genel & Evrak İşleri",
        "cat2": "IE 300 & IE 400 Detayları",
        "q1_lbl": "Staj yeri nasıl ayarlanır?",
        "q1_prmpt": "Yaz stajı organizasyonunu nasıl ayarlayabilirim?",
        "q2_lbl": "Gerekli evraklar nelerdir?",
        "q2_prmpt": "Yaz stajına resmi olarak kayıt olmak için gerekli evraklar nelerdir?",
        "q3_lbl": "Hangi departmanları ziyaret etmeliyim?",
        "q3_prmpt": "Yaz stajım sırasında organizasyonun hangi departmanlarını ziyaret etmeliyim?",
        "q4_lbl": "Bankalar IE 400 için uygun mu?",
        "q4_prmpt": "Bir bankadaki hangi tür şubeler veya bölümler IE 400 staj organizasyonu için uygundur?",
        "q5_lbl": "Hizmet firmalarında üretim soruları?",
        "q5_prmpt": "Staj hizmet organizasyonunda yapılıyorsa, öğrenci üretim odaklı alanlardaki soruları nasıl yanıtlamalıdır?",
        "q6_lbl": "Yeterli problem tanımı (IE400)?",
        "q6_prmpt": "IE 400 raporlarında 'yeterli' bir Endüstri Mühendisliği problem tanımı neleri içerir?",
        "q7_lbl": "Bunun yerine tam zamanlı proje yapabilir miyim?",
        "q7_prmpt": "Yaz stajı organizasyonunda tam zamanlı bir projede çalışmayı planlıyorsam, proje çalışmasını sunmam mümkün mü?",
        "q8_lbl": "Aynı firmadaki birden fazla öğrenci için kurallar?",
        "q8_prmpt": "Aynı organizasyonda birkaç öğrenci staj yaparsa ne tür bir raporlama kabul edilir?",
        "placeholder": "Örn., IE 300 için gereksinimler nelerdir?",
        "thinking": "Resmi yönergeler inceleniyor...",
        "system_instruction_lang": "You must respond entirely in Turkish. (Kullanıcıya tamamen Türkçe cevap vermelisin.)"
    }
}

# --- SIDEBAR & LANGUAGE SELECTION ---
with st.sidebar:
    selected_lang = st.radio("Language / Dil", ["ENG", "TR"], horizontal=True, label_visibility="collapsed")
    st.markdown("---")
    try:
        st.image("metu_logo.png", use_container_width=True)
    except:
        st.error("Logo missing. Add 'metu_logo.png' to folder.")
        
t = lang_dict[selected_lang]

# --- REUSABLE UI HEADER FUNCTION ---
def render_metu_header():
    st.markdown(f"<div class='header-text'>{t['title']}</div>", unsafe_allow_html=True)
    st.markdown(f"""
        <div class='metu-topbar'>{t['topbar']}</div>
        <div class='metu-navbar'>{t['subtitle']}</div>
    """, unsafe_allow_html=True)

# --- LOCAL KNOWLEDGE BASE LOADER ---
@st.cache_data 
def load_knowledge_base():
    kb_text = ""
    kb_folder = "knowledge_base"

    if not os.path.exists(kb_folder):
        return "Warning: 'knowledge_base' folder not found. Please create it and add your files."

    for filename in os.listdir(kb_folder):
        filepath = os.path.join(kb_folder, filename)
        try:
            if filename.endswith(".txt"):
                with open(filepath, "r", encoding="utf-8") as f:
                    kb_text += f"\n\n--- Source: {filename} ---\n" + f.read()
            elif filename.endswith(".pdf"):
                reader = PdfReader(filepath)
                kb_text += f"\n\n--- Source: {filename} ---\n"
                for page in reader.pages:
                    text = page.extract_text()
                    if text: kb_text += text + "\n"
            elif filename.endswith(".pptx"):
                prs = Presentation(filepath)
                kb_text += f"\n\n--- Source: {filename} ---\n"
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            kb_text += shape.text + "\n"
            elif filename.endswith(".docx"):
                doc = Document(filepath)
                kb_text += f"\n\n--- Source: {filename} ---\n"
                for para in doc.paragraphs:
                    if para.text.strip():
                        kb_text += para.text + "\n"
        except Exception as e:
            print(f"Error reading {filename}: {e}")

    return kb_text

# --- INITIALIZE OPENAI SECURELY ---
try:
    client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])
except KeyError:
    st.error("API Key not found! Please configure your Streamlit Secrets.")
    st.stop()

# --- SESSION STATE INITIALIZATION ---
if "messages" not in st.session_state:
    st.session_state.messages = []
    
if "quick_prompt" not in st.session_state:
    st.session_state.quick_prompt = None

# --- SYSTEM INSTRUCTIONS ---
kb_text = load_knowledge_base()

system_instruction = f"""
You are an intelligent virtual consultant specifically designed to assist 
Middle East Technical University (METU) Industrial Engineering students 
with their Summer Practices (IE 300 and IE 400). 

Below is the official documentation. You MUST use ONLY this information to answer user queries. 
Do not make up deadlines or rules.

<official_documentation>
{kb_text}
</official_documentation>

CRITICAL INSTRUCTION: If a user asks a question that falls outside the scope 
of METU-IE Summer Practice (e.g., general programming, other universities, or unrelated topics), 
you must politely decline to answer and redirect them to ask about IE 300 or IE 400 summer practices.

{t['system_instruction_lang']}
"""

# --- RENDER MAIN UI ---
render_metu_header()

# --- CONTINUE ENHANCED SIDEBAR ---
with st.sidebar:
    st.markdown("---")
    
    # Quick Links Section
    st.header(t["links_header"])
    st.markdown(f"""
        <a href="https://ie.metu.edu.tr/en/current-semester" target="_blank" class="sidebar-link">📚 {t['link_curr']}</a>
        <a href="https://sp-ie.metu.edu.tr/en" target="_blank" class="sidebar-link">🌐 {t['link_sp']}</a>
        <a href="https://ie.metu.edu.tr/en" target="_blank" class="sidebar-link">🏠 {t['link_main']}</a>
    """, unsafe_allow_html=True)
    
    st.markdown("---")
    
    # FAQ Section
    st.header(t["faq_header"])
    
    st.subheader(t["cat1"])
    if st.button(t["q1_lbl"]):
        st.session_state.quick_prompt = t["q1_prmpt"]
    if st.button(t["q2_lbl"]):
        st.session_state.quick_prompt = t["q2_prmpt"]
    if st.button(t["q3_lbl"]):
        st.session_state.quick_prompt = t["q3_prmpt"]
        
    st.subheader(t["cat2"])
    if st.button(t["q4_lbl"]):
        st.session_state.quick_prompt = t["q4_prmpt"]
    if st.button(t["q5_lbl"]):
        st.session_state.quick_prompt = t["q5_prmpt"]
    if st.button(t["q6_lbl"]):
        st.session_state.quick_prompt = t["q6_prmpt"]
    if st.button(t["q7_lbl"]):
        st.session_state.quick_prompt = t["q7_prmpt"]
    if st.button(t["q8_lbl"]):
        st.session_state.quick_prompt = t["q8_prmpt"]

    st.markdown("---")
    
    # Contact Information at the bottom
    st.header(t["contact_hdr"])
    st.markdown("""
        **Tel:** +90 (312) 210 4796  
        **Fax:** +90 (312) 210 4786  
        **E-Mail:** [ie-staj@metu.edu.tr](mailto:ie-staj@metu.edu.tr)  
        **Address:** METU Industrial Engineering, 06800
    """)

# --- CHAT RENDERING & LOGIC ---
for msg in st.session_state.messages:
    st.chat_message(msg["role"]).write(msg["content"])

prompt = st.chat_input(t["placeholder"])

if st.session_state.quick_prompt:
    prompt = st.session_state.quick_prompt
    st.session_state.quick_prompt = None

if prompt:
    st.session_state.messages.append({"role": "user", "content": prompt})
    st.chat_message("user").write(prompt)

    # Wrap the API call in a spinner to show the thinking animation
    with st.spinner(t["thinking"]):
        try:
            api_messages = [{"role": "system", "content": system_instruction}] + st.session_state.messages

            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=api_messages,
                temperature=0.2
            )
            
            bot_reply = response.choices[0].message.content
            
            st.session_state.messages.append({"role": "assistant", "content": bot_reply})
            st.chat_message("assistant").write(bot_reply)
            
        except Exception as e:
            st.error(f"An error occurred: {e}")